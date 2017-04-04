from pptx import Presentation # pip install python-pptx
from pptx.util import Inches
from os import listdir
from os.path import isfile, join

def Output_to_PPT(filename):
    # filename = "GBPUSD"
    # modify the path by changing the file name
    search_dir ='C:\\Users\\YYJ\\Desktop\\FIN580\\Homework1\\VolatilityForecasting\\src\\'+filename+'\\'
    os.chdir(search_dir)
    files = filter(os.path.isfile, os.listdir(search_dir))
    files = [os.path.join(search_dir, f) for f in files]  # add path to each file
    files.sort(key=lambda x: os.path.getmtime(x))

    # onlyfiles = [f for f in listdir(mypath) if isfile(join(mypath, f))]
    # onlyfiles.sort(key=lambda x: os.path.getmtime(x))

    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]

    for i in range(len(files)):
        img_path = files[i]
        slide = prs.slides.add_slide(blank_slide_layout)
        left = Inches(0.8)
        height = Inches(5.5)
        pic = slide.shapes.add_picture(img_path, left, top, height=height)

    prs.save(filename+'.pptx')