from pptx import Presentation # pip install python-pptx
from pptx.util import Inches
from os import listdir
from os.path import isfile, join

def Output_to_PPT(filename):
    # filename = "GBPUSD"
    # modify the path by changing the file name
    search_dir ='C:\\Users\\YYJ\\Desktop\\FIN580\\Homework1\\VolatilityForecasting\\src\\'+filename+'\\'
    os.chdir(search_dir)

    files_all_type = filter(os.path.isfile, os.listdir(search_dir))
    files_specific_type=[]
    for file in files_all_type:
        if file.endswith(".png"):
            files_specific_type.append(os.path.join(search_dir, file))  # add path to each file
        files_specific_type.sort(key=lambda x: os.path.getmtime(x))

    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]

    for i in range(len(files_specific_type)):
        img_path = files_specific_type[i]
        slide = prs.slides.add_slide(blank_slide_layout)
        left = Inches(0.8)
        height = Inches(5.5)
        pic = slide.shapes.add_picture(img_path, left, top, height=height)

    prs.save(filename+'.pptx')